"""
Build presentation.pptx for CS495 Capstone — Prompt Injection Robustness
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

ROOT = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(ROOT, "figures")

# ── Colour palette ─────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x11, 0x17)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF0, 0xF2, 0xF5)
ACCENT_BLUE = RGBColor(0x29, 0x80, 0xB9)
ACCENT_GOLD = RGBColor(0xF3, 0x9C, 0x12)
ACCENT_RED  = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_GRN  = RGBColor(0x27, 0xAE, 0x60)
MID_GREY    = RGBColor(0x55, 0x65, 0x75)
DIM_WHITE   = RGBColor(0xCC, 0xD1, 0xD9)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────
def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=Pt(18),
                bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True,
                italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_para(tf, text, font_size=Pt(16), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(4), bullet_char=None,
             italic=False, indent_level=0):
    from pptx.util import Pt as pt2
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.level = indent_level
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def slide_header(slide, title, subtitle=None, dark=True):
    """Adds a top header bar with title (and optional subtitle)."""
    bg_col = NAVY if dark else LIGHT_GREY
    text_col = WHITE if dark else NAVY

    # header bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), fill_color=ACCENT_BLUE)

    add_textbox(slide, Inches(0.35), Inches(0.12), Inches(12.3), Inches(0.75),
                title, font_size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    if subtitle:
        add_textbox(slide, Inches(0.35), Inches(0.82), Inches(12.3), Inches(0.38),
                    subtitle, font_size=Pt(14), bold=False, color=DIM_WHITE, align=PP_ALIGN.LEFT)


def footer_bar(slide, text="Kelan Huang  ·  CS495 Capstone  ·  Bellevue College  ·  Spring 2026"):
    add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), fill_color=MID_GREY)
    add_textbox(slide, Inches(0.3), Inches(7.17), Inches(12.7), Inches(0.3),
                text, font_size=Pt(9), bold=False, color=DIM_WHITE,
                align=PP_ALIGN.LEFT)


def add_image(slide, path, left, top, width, height=None):
    if not os.path.exists(path):
        return None
    if height:
        return slide.shapes.add_picture(path, left, top, width, height)
    else:
        return slide.shapes.add_picture(path, left, top, width)


def bullet_box(slide, left, top, width, height, items,
               font_size=Pt(16), color=WHITE, title=None, title_color=ACCENT_GOLD,
               title_size=Pt(17), indent_items=None):
    """Adds a textbox with optional title + bulleted items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.bold = True
        run.font.size = title_size
        run.font.color.rgb = title_color
        run.font.name = "Calibri"

    for i, item in enumerate(items):
        indent = indent_items[i] if indent_items else 0
        p = tf.add_paragraph() if (not first or title) else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.level = indent
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = ("    " * indent + "• " if indent > 0 else "• ") + item
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
# BUILD SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def build_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── SLIDE 1 — Title ────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    # Decorative accent stripe
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill_color=ACCENT_BLUE)
    add_rect(sl, Inches(0.18), Inches(0), Inches(0.05), SLIDE_H, fill_color=ACCENT_GOLD)

    add_textbox(sl, Inches(0.55), Inches(1.6), Inches(11.8), Inches(1.8),
                "Comparing Prompt Injection\nRobustness Across Pre-Trained LLMs",
                font_size=Pt(38), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_textbox(sl, Inches(0.55), Inches(3.55), Inches(11.0), Inches(0.55),
                "CS495 Capstone Research  ·  Spring 2026",
                font_size=Pt(20), bold=False, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)

    add_textbox(sl, Inches(0.55), Inches(4.2), Inches(11.0), Inches(0.45),
                "Kelan Huang  |  Bellevue College",
                font_size=Pt(17), bold=False, color=DIM_WHITE, align=PP_ALIGN.LEFT)

    # Model list strip
    add_rect(sl, Inches(0.55), Inches(5.3), Inches(11.7), Inches(0.05), fill_color=ACCENT_BLUE)
    add_textbox(sl, Inches(0.55), Inches(5.45), Inches(11.7), Inches(0.5),
                "Llama 3.1 8B  ·  GPT-5.5  ·  Gemini 3 Flash  ·  Qwen 3.6 35B  ·  DeepSeek V4 Pro  ·  Claude Sonnet 4.6",
                font_size=Pt(13), bold=False, color=DIM_WHITE, align=PP_ALIGN.LEFT)

    footer_bar(sl)

    # ── SLIDE 2 — Introduction ─────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, "Introduction", "Why does prompt injection matter?")
    footer_bar(sl)

    # Left column
    bullet_box(sl, Inches(0.35), Inches(1.3), Inches(6.2), Inches(5.5),
        title="Security Context",
        title_color=ACCENT_GOLD,
        title_size=Pt(17),
        items=[
            "Prompt injection = top threat to LLM-powered apps (OWASP 2024)",
            "Commercial models: 80–100% ASR in targeted scenarios",
            "Orgs deploying smaller, self-hosted models lack robustness data",
        ],
        font_size=Pt(15), color=WHITE)

    bullet_box(sl, Inches(6.8), Inches(1.3), Inches(6.15), Inches(5.5),
        title="Research Questions",
        title_color=ACCENT_GOLD,
        title_size=Pt(17),
        items=[
            "Which models are most/least vulnerable to baseline attacks?",
            "Does switching attack language (Mandarin/Swahili/Welsh) raise ASR?",
            "Can an adversarial LLM auto-generate effective attacks?",
        ],
        font_size=Pt(15), color=WHITE)

    # Scope strip
    add_rect(sl, Inches(0.35), Inches(5.4), Inches(12.6), Inches(1.55),
             fill_color=RGBColor(0x1A, 0x24, 0x33))
    add_textbox(sl, Inches(0.55), Inches(5.45), Inches(12.2), Inches(0.4),
                "Scope",
                font_size=Pt(15), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.LEFT)
    add_textbox(sl, Inches(0.55), Inches(5.8), Inches(12.2), Inches(0.9),
                "6 LLMs  ·  10 attack techniques  ·  2 agent domains (Cooking, Health)  ·  4 languages (English, Mandarin, Swahili, Welsh)  ·  Runs on consumer gaming GPU",
                font_size=Pt(14), bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    # ── SLIDE 3 — Literature Review ────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, "Literature Review", "What has prior research found?")
    footer_bar(sl)

    # Three columns
    col_w = Inches(3.9)
    for ci, (title, bullets, col) in enumerate([
        ("Attack Taxonomy", [
            "Prompt hijacking — subvert model intent",
            "Prompt extraction — leak system instructions",
            "HackAPrompt & Tensor Trust: 700K+ adversarial examples",
            "Reasoning models more robust overall, but vulnerable to logic attacks (TAP)",
        ], ACCENT_BLUE),
        ("Attack Vectors", [
            "Obfuscation: Base64, ROT13 encoding",
            "Distraction & fake completion injections",
            "Tool hijacking in agentic pipelines",
            "Social engineering (authority, roleplay)",
        ], ACCENT_GOLD),
        ("Defence & Risks", [
            "StruQ (structured queries) & SecAlign",
            "Attention tracking detects distraction layers",
            "Medical LLM attacks: 94.4% unsafe recommendation rate",
            "Non-English safety gaps; overdefence trade-off",
        ], ACCENT_RED),
    ]):
        left = Inches(0.35 + ci * 4.3)
        add_rect(sl, left, Inches(1.3), col_w, Inches(0.38), fill_color=col)
        add_textbox(sl, left + Inches(0.12), Inches(1.32), col_w - Inches(0.24), Inches(0.36),
                    title, font_size=Pt(16), bold=True, color=WHITE)
        bullet_box(sl, left, Inches(1.7), col_w, Inches(5.2),
                   items=bullets, font_size=Pt(14), color=WHITE)

    # ── SLIDE 4 — Methodology ──────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, "Methodology", "Experimental design overview")
    footer_bar(sl)

    # Left column — models & phases
    bullet_box(sl, Inches(0.35), Inches(1.3), Inches(5.8), Inches(5.6),
        title="Models & Domains",
        title_color=ACCENT_BLUE,
        items=[
            "Claude Sonnet 4.6  (commercial)",
            "GPT-5.5  (commercial)",
            "Gemini 3 Flash  (commercial)",
            "Llama 3.1 8B  (open-source)",
            "Qwen 3.6 35B A3B MTP  (open-source)",
            "DeepSeek V4 Pro  (open-source)",
            "Domains: Cooking Assistant · Health Assistant",
        ],
        font_size=Pt(15), color=WHITE)

    # Right column — phases
    bullet_box(sl, Inches(6.4), Inches(1.3), Inches(6.55), Inches(5.6),
        title="Experimental Phases",
        title_color=ACCENT_GOLD,
        items=[
            "Phase I — 5 baseline attacks × 2 domains × 5 reps",
            "  Naive · Roleplay/DAN · Fake Completion · Extraction · Base64",
            "Phase IIB — Top-3 attacks in Mandarin, Swahili, Welsh",
            "Phase IIC — Llama 3.1 8B generates payloads vs Claude",
            "Phase III — Qwen family: 9B vs 27B vs 35B comparison",
            "Metric: ASR = successes ÷ non-ambiguous runs",
            "Error bars: ±2 SEM (≈95% CI)",
            "Scorer: automated LLM rubric (rubric v2)",
        ],
        indent_items=[0, 1, 0, 0, 0, 0, 0, 0],
        font_size=Pt(14), color=WHITE)

    # ── SLIDE 5 — Results Part 1 (Phase I) ────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, LIGHT_GREY)
    slide_header(sl, "Results — Phase I: Baseline Attacks",
                 "Llama 3.1 8B most vulnerable · Claude Sonnet 4.6 & DeepSeek most robust")
    footer_bar(sl, "")

    # Fig 1 (ASR by model)
    add_image(sl, os.path.join(FIG, "fig1_p1_asr_by_model.png"),
              Inches(0.3), Inches(1.25), Inches(6.3))

    # Fig 2 (ASR by attack)
    add_image(sl, os.path.join(FIG, "fig2_p1_asr_by_attack.png"),
              Inches(6.75), Inches(1.25), Inches(6.25))

    # Key insight ribbon
    add_rect(sl, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.5),
             fill_color=RGBColor(0xE8, 0xF4, 0xFD))
    add_textbox(sl, Inches(0.45), Inches(6.57), Inches(12.4), Inches(0.45),
                "Key insight:  Health domain ASR (59.7%) >> Cooking (18.4%) — domain context is a strong vulnerability moderator",
                font_size=Pt(13), bold=True,
                color=RGBColor(0x0D, 0x11, 0x17), align=PP_ALIGN.LEFT)

    # ── SLIDE 6 — Results Part 2 (Phase IIB + IIC + III) ──────────────────
    sl = blank_slide(prs)
    fill_bg(sl, LIGHT_GREY)
    slide_header(sl, "Results — Phase IIB · IIC · III",
                 "Multilingual & adversarial attacks · Qwen family comparison")
    footer_bar(sl, "")

    # Fig 4 (language ASR) — top-left
    add_image(sl, os.path.join(FIG, "fig4_p2b_asr_by_language.png"),
              Inches(0.3), Inches(1.25), Inches(4.2))

    # Fig 6 (IIC comparison) — top-middle
    add_image(sl, os.path.join(FIG, "fig6_p2c_asr_comparison.png"),
              Inches(4.65), Inches(1.25), Inches(4.2))

    # Fig 14 (Phase III Qwen) — top-right
    add_image(sl, os.path.join(FIG, "fig14_p3_phase1_asr.png"),
              Inches(9.0), Inches(1.25), Inches(4.05))

    # Three caption bars
    for left, w, label, col in [
        (Inches(0.3),   Inches(4.2),  "IIB: Language switching ≠ bypass",   ACCENT_BLUE),
        (Inches(4.65),  Inches(4.2),  "IIC: LLM-generated payloads (n=15)", ACCENT_RED),
        (Inches(9.0),   Inches(4.05), "III: Qwen 9B → 35B, similar ASR",   ACCENT_GRN),
    ]:
        add_rect(sl, left, Inches(5.55), w, Inches(0.32), fill_color=col)
        add_textbox(sl, left + Inches(0.08), Inches(5.57), w - Inches(0.12), Inches(0.3),
                    label, font_size=Pt(12), bold=True, color=WHITE)

    # Summary bullets
    add_textbox(sl, Inches(0.35), Inches(5.95), Inches(12.6), Inches(1.1),
                "• Mandarin/Swahili/Welsh ASR (43–47%) not significantly higher than English baseline (43.2%)   "
                "• Llama-generated attacks: 33% overall ASR vs Claude (limited n=15)   "
                "• Qwen scaling: health-domain ASR rises 62% → 80% from 9B → 35B — larger ≠ safer",
                font_size=Pt(13), bold=False,
                color=RGBColor(0x1A, 0x1A, 0x2E), align=PP_ALIGN.LEFT)

    # ── SLIDE 7 — Limitations & Future Work ───────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, "Limitations & Future Work")
    footer_bar(sl)

    bullet_box(sl, Inches(0.35), Inches(1.3), Inches(6.0), Inches(5.6),
        title="Limitations",
        title_color=ACCENT_RED,
        items=[
            "Phase IIC: only 5 reps/attack — Naive 100% likely a small-sample artifact",
            "Automated rubric may misclassify subtle responses (AMBIGUOUS category)",
            "Static hand-crafted prompts; adaptive attacker could be more adversarial",
            "LM Studio inference variability not fully controlled",
        ],
        font_size=Pt(15), color=WHITE)

    bullet_box(sl, Inches(6.65), Inches(1.3), Inches(6.3), Inches(5.6),
        title="Next Steps",
        title_color=ACCENT_GRN,
        items=[
            "Phase IV: LoRA fine-tuning on Llama 3.1 8B to reduce vulnerability",
            "Expand Phase IIC to ≥30 reps for statistical power",
            "Gemma 4 dense (31B) vs Gemma 4 MoE (26B A4B) comparison",
            "Dense vs MoE security trade-off across Alibaba Qwen & Google Gemma",
        ],
        font_size=Pt(15), color=WHITE)

    # ── SLIDE 8 — Key Findings / Conclusion ───────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, "Key Findings", "What this research tells us")
    footer_bar(sl)

    findings = [
        ("1", "Model vulnerability spans 10–80% ASR", "Llama 3.1 8B (79.6%) most susceptible; Claude (10%) & DeepSeek (12%) most robust", ACCENT_RED),
        ("2", "Language switching ≠ bypass", "Mandarin/Swahili/Welsh ASR comparable to English; no systematic amplification", ACCENT_BLUE),
        ("3", "LLM-generated attacks fall short", "Llama-generated payloads fail to systematically beat a strong defender (Claude)", ACCENT_GOLD),
        ("4", "Domain context is a key moderator", "Health agent: 59.7% ASR vs Cooking: 18.4% — same model, very different risk", ACCENT_GRN),
        ("5", "Larger Qwen ≠ safer", "Qwen 35B health ASR (80%) exceeds 9B (62%) — scaling up can increase vulnerability", RGBColor(0x8E, 0x44, 0xAD)),
    ]

    row_h = Inches(0.95)
    for i, (num, headline, detail, col) in enumerate(findings):
        top = Inches(1.3) + i * row_h
        add_rect(sl, Inches(0.35), top, Inches(0.5), row_h - Inches(0.06),
                 fill_color=col)
        add_textbox(sl, Inches(0.35), top + Inches(0.18), Inches(0.5), Inches(0.5),
                    num, font_size=Pt(22), bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, Inches(0.95), top + Inches(0.04), Inches(11.9), Inches(0.38),
                    headline, font_size=Pt(16), bold=True, color=WHITE)
        add_textbox(sl, Inches(0.95), top + Inches(0.42), Inches(11.9), Inches(0.44),
                    detail, font_size=Pt(13), bold=False, color=DIM_WHITE)

    out_path = os.path.join(ROOT, "presentation.pptx")
    prs.save(out_path)
    print(f"Saved → {out_path}")
    return out_path


if __name__ == "__main__":
    build_prs()
